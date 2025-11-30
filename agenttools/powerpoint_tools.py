"""PowerPoint presentation manipulation tools using python-pptx."""

from typing import Annotated
from semantic_kernel.functions import kernel_function

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class PowerPointTools:
    """Tools for creating and editing PowerPoint presentations."""

    def __init__(self):
        """Initialize PowerPointTools."""
        if not PPTX_AVAILABLE:
            print("Warning: python-pptx not installed. Install with: pip install python-pptx")

    @kernel_function(
        name="create_presentation",
        description="Create a new PowerPoint presentation at the specified path"
    )
    def create_presentation(
        self,
        file_path: Annotated[str, "Full path where the presentation will be created"]
    ) -> Annotated[str, "Success or error message"]:
        """Create a new blank PowerPoint presentation."""
        if not PPTX_AVAILABLE:
            return "Error: python-pptx is not installed"
        
        try:
            prs = Presentation()
            prs.save(file_path)
            return f"PowerPoint presentation created successfully at: {file_path}"
        except Exception as err:
            return f"Error creating presentation: {str(err)}"

    @kernel_function(
        name="add_title_slide",
        description="Add a title slide to a PowerPoint presentation"
    )
    def add_title_slide(
        self,
        file_path: Annotated[str, "Path to the PowerPoint presentation"],
        title: Annotated[str, "Slide title text"],
        subtitle: Annotated[str, "Slide subtitle text"] = ""
    ) -> Annotated[str, "Success or error message"]:
        """Add a title slide to an existing presentation."""
        if not PPTX_AVAILABLE:
            return "Error: python-pptx is not installed"
        
        try:
            prs = Presentation(file_path)
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            
            slide.shapes.title.text = title
            if subtitle and len(slide.placeholders) > 1:
                slide.placeholders[1].text = subtitle
            
            prs.save(file_path)
            return f"Title slide added successfully to: {file_path}"
        except Exception as err:
            return f"Error adding title slide: {str(err)}"

    @kernel_function(
        name="add_content_slide",
        description="Add a content slide with title and bullet points"
    )
    def add_content_slide(
        self,
        file_path: Annotated[str, "Path to the PowerPoint presentation"],
        title: Annotated[str, "Slide title"],
        content: Annotated[str, "Comma-separated bullet points"]
    ) -> Annotated[str, "Success or error message"]:
        """Add a content slide with bullet points."""
        if not PPTX_AVAILABLE:
            return "Error: python-pptx is not installed"
        
        try:
            prs = Presentation(file_path)
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            
            slide.shapes.title.text = title
            
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            
            bullets = [item.strip() for item in content.split(',')]
            for bullet in bullets:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0
            
            prs.save(file_path)
            return f"Content slide with {len(bullets)} bullets added successfully"
        except Exception as err:
            return f"Error adding content slide: {str(err)}"

    @kernel_function(
        name="add_blank_slide",
        description="Add a blank slide to a PowerPoint presentation"
    )
    def add_blank_slide(
        self,
        file_path: Annotated[str, "Path to the PowerPoint presentation"]
    ) -> Annotated[str, "Success or error message"]:
        """Add a blank slide to an existing presentation."""
        if not PPTX_AVAILABLE:
            return "Error: python-pptx is not installed"
        
        try:
            prs = Presentation(file_path)
            blank_slide_layout = prs.slide_layouts[6]
            prs.slides.add_slide(blank_slide_layout)
            prs.save(file_path)
            return f"Blank slide added successfully to: {file_path}"
        except Exception as err:
            return f"Error adding blank slide: {str(err)}"

    @kernel_function(
        name="add_text_box",
        description="Add a text box to a specific slide"
    )
    def add_text_box(
        self,
        file_path: Annotated[str, "Path to the PowerPoint presentation"],
        slide_index: Annotated[int, "Index of the slide (0-based)"],
        text: Annotated[str, "Text content for the text box"],
        left: Annotated[float, "Left position in inches"] = 1.0,
        top: Annotated[float, "Top position in inches"] = 1.0,
        width: Annotated[float, "Width in inches"] = 5.0,
        height: Annotated[float, "Height in inches"] = 1.0
    ) -> Annotated[str, "Success or error message"]:
        """Add a text box to a specific slide."""
        if not PPTX_AVAILABLE:
            return "Error: python-pptx is not installed"
        
        try:
            prs = Presentation(file_path)
            if slide_index >= len(prs.slides):
                return f"Error: Slide index {slide_index} out of range"
            
            slide = prs.slides[slide_index]
            textbox = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            text_frame = textbox.text_frame
            text_frame.text = text
            
            prs.save(file_path)
            return f"Text box added successfully to slide {slide_index}"
        except Exception as err:
            return f"Error adding text box: {str(err)}"

    @kernel_function(
        name="get_slide_count",
        description="Get the number of slides in a PowerPoint presentation"
    )
    def get_slide_count(
        self,
        file_path: Annotated[str, "Path to the PowerPoint presentation"]
    ) -> Annotated[str, "Number of slides or error message"]:
        """Get the total number of slides in a presentation."""
        if not PPTX_AVAILABLE:
            return "Error: python-pptx is not installed"
        
        try:
            prs = Presentation(file_path)
            return f"Presentation has {len(prs.slides)} slides"
        except Exception as err:
            return f"Error reading presentation: {str(err)}"

    @kernel_function(
        name="read_slide_text",
        description="Read all text from a specific slide"
    )
    def read_slide_text(
        self,
        file_path: Annotated[str, "Path to the PowerPoint presentation"],
        slide_index: Annotated[int, "Index of the slide (0-based)"]
    ) -> Annotated[str, "Slide text content or error message"]:
        """Read all text from a specific slide."""
        if not PPTX_AVAILABLE:
            return "Error: python-pptx is not installed"
        
        try:
            prs = Presentation(file_path)
            if slide_index >= len(prs.slides):
                return f"Error: Slide index {slide_index} out of range"
            
            slide = prs.slides[slide_index]
            text_content = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
            
            return "\n".join(text_content) if text_content else "No text found on slide"
        except Exception as err:
            return f"Error reading slide: {str(err)}"
